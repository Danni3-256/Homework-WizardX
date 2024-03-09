import os
import io
import openai
import docx2txt
import PyPDF2
from pptx import Presentation

openai.api_key = os.environ.get("OPENAI_API_KEY")

# Loading Files
def load_text_from_docx(file_content):
    file_io = io.BytesIO(file_content)
    text = docx2txt.process(file_io)
    return text

def load_text_from_pdf(file_path):
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfFileReader(file)
        text = ''
        for page_num in range(reader.numPages):
            text += reader.getPage(page_num).extractText()
    return text

def load_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def load_student_homework(uploaded_file):
    if uploaded_file is None:
        raise ValueError ('No File Uploaded')
    
    _, file_extension = os.path.splitext(uploaded_file.name.lower())

    content = uploaded_file.read()

    if file_extension == '.docx':
        return load_text_from_docx(content)
    elif file_extension == '.pdf':
        return load_text_from_pdf(content)
    elif file_extension == '.pptx':
        return load_text_from_pptx(content)
    else:
        raise ValueError("Unsupported file format")
    
# Generating Feed Back using chatGPT

def generate_feedback(prompt_template, homework_text):
    prompt = prompt_template.format(homework_text)
    response = openai.Completion.create(
        model="gpt-3.5-turbo-instruct", 
        prompt=prompt,
        temperature=0.7,
        max_tokens=500,
        n=1,
        stop=None,
    )
    feedback = response.choices[0].text.strip()
    return feedback

# Prompting the Model

def create_prompt_template():
    prompt_template = "First, extract student and course details on a new line and then provide feedback on the following assignment and assign marks out of 20.:\n\n{}\n\nFeedback:"
    return prompt_template


def main():
    homework_file_path = ""
    student_homework = load_student_homework(homework_file_path)
    prompt_template = create_prompt_template()
    feedback = generate_feedback(prompt_template, student_homework)

    print("Generated Feedback:")
    print(feedback)

if __name__ == "__main__":
    main()
